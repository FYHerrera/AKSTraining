"""
Generate PPT presentations for all 10 AKS course lessons.
Creates English and Spanish versions for each lesson.

Usage:
    python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────────────
AZURE_BLUE    = RGBColor(0, 120, 212)
DARK_BLUE     = RGBColor(0, 78, 140)
LIGHT_BLUE    = RGBColor(0, 153, 255)
WHITE         = RGBColor(255, 255, 255)
BLACK         = RGBColor(33, 33, 33)
GRAY          = RGBColor(100, 100, 100)
LIGHT_GRAY    = RGBColor(200, 200, 200)
GREEN         = RGBColor(16, 137, 62)
RED           = RGBColor(209, 52, 56)
ORANGE        = RGBColor(255, 140, 0)
CODE_BG       = RGBColor(40, 44, 52)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
LESSONS_DIR = os.path.join(BASE_DIR, "lecciones")


def add_bg_shape(slide, color=AZURE_BLUE):
    """Add a colored bar at the top of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bottom_bar(slide, text="AKS Troubleshooting Course"):
    """Add a subtle bottom bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.05), Inches(13.33), Inches(0.45)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs, title, subtitle, lesson_num):
    """Add a title slide with Azure branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full blue background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = AZURE_BLUE
    bg.line.fill.background()

    # Accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3.2), Inches(13.33), Inches(0.06)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = LIGHT_BLUE
    stripe.line.fill.background()

    # Lesson number
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Lesson {lesson_num:02d}" if lesson_num else ""
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.font.bold = False

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 220, 240)

    # Footer
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "AKS Troubleshooting Course | Azure Kubernetes Service"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(150, 190, 230)


def add_section_slide(prs, title, subtitle=""):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = LIGHT_BLUE


def add_content_slide(prs, title, bullets, footer="AKS Troubleshooting Course"):
    """Add a content slide with bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_bottom_bar(slide, footer)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.3))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        if bullet.startswith(">>"):
            p.text = bullet[2:].strip()
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY
        else:
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK

        p.space_after = Pt(6)


def add_code_slide(prs, title, code_lines, footer="AKS Troubleshooting Course"):
    """Add a slide with a code block."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_bottom_bar(slide, footer)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Code block background
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(12), Inches(5.2)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = CODE_BG
    code_box.line.color.rgb = RGBColor(60, 63, 70)

    # Code text
    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.4), Inches(4.9))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(171, 178, 191)
        p.font.name = "Consolas"
        p.space_after = Pt(2)


def add_table_slide(prs, title, headers, rows, footer="AKS Troubleshooting Course"):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide)
    add_bottom_bar(slide, footer)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1
    tbl = slide.shapes.add_table(
        table_rows, cols, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5 * table_rows)
    ).table

    # Header
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZURE_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(240, 245, 250)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = BLACK


def add_lab_slide(prs, lab_num, title, description, command, lang="en"):
    """Add a lab exercise slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(25, 35, 55)
    bg.line.fill.background()

    # Lab badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(2), Inches(0.6)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GREEN
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = f"Lab {lab_num:02d}" if lang == "en" else f"Lab {lab_num:02d}"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Description
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(180, 200, 220)

    # Command box
    cmd_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11), Inches(1.4)
    )
    cmd_bg.fill.solid()
    cmd_bg.fill.fore_color.rgb = CODE_BG
    cmd_bg.line.color.rgb = GREEN

    txBox = slide.shapes.add_textbox(Inches(1.1), Inches(5.15), Inches(10.5), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"$ chmod +x {command}"
    p.font.size = Pt(18)
    p.font.color.rgb = GREEN
    p.font.name = "Consolas"
    p2 = tf.add_paragraph()
    p2.text = f"$ ./{command}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = GREEN
    p2.font.name = "Consolas"


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


# ═══════════════════════════════════════════════════════════════════════════
# LESSON DEFINITIONS - Both languages
# ═══════════════════════════════════════════════════════════════════════════

LESSONS = {
    1: {
        "en": {
            "title": "kubectl Fundamentals & AKS Architecture",
            "subtitle": "Learn the essential commands and understand how AKS clusters work",
            "slides": [
                ("section", "AKS Architecture", "Understanding the building blocks"),
                ("content", "Control Plane (Managed by Azure)", [
                    "API Server – Entry point for all kubectl operations",
                    "etcd – Key-value store for cluster state",
                    "Scheduler – Assigns pods to nodes",
                    "Controller Manager – Maintains desired state",
                    "Cloud Controller Manager – Azure integration",
                ]),
                ("content", "Worker Nodes", [
                    "kubelet – Agent that runs pods on each node",
                    "kube-proxy – Network rules on each node",
                    "Container Runtime – Runs the actual containers",
                    "CoreDNS – Internal DNS resolution",
                ]),
                ("section", "Connecting to the Cluster", ""),
                ("code", "Connect to AKS", [
                    "# Login to Azure",
                    "az login",
                    "",
                    "# Get cluster credentials",
                    "az aks get-credentials --resource-group <rg> --name <cluster>",
                    "",
                    "# Verify connection",
                    "kubectl cluster-info",
                ]),
                ("section", "Essential kubectl Commands", ""),
                ("code", "Viewing Resources", [
                    "# Nodes",
                    "kubectl get nodes -o wide",
                    "",
                    "# Pods (all namespaces)",
                    "kubectl get pods -A",
                    "",
                    "# Everything",
                    "kubectl get all -n <namespace>",
                ]),
                ("code", "Diagnosing Resources", [
                    "# Describe shows EVENTS - key for troubleshooting",
                    "kubectl describe pod <name>",
                    "kubectl describe node <name>",
                    "",
                    "# Logs",
                    "kubectl logs <pod-name>",
                    "kubectl logs <pod-name> --previous",
                    "",
                    "# Execute inside a pod",
                    "kubectl exec -it <pod> -- /bin/sh",
                ]),
                ("content", "Namespaces", [
                    "default – Resources go here if no namespace specified",
                    "kube-system – System components (CoreDNS, kube-proxy, etc.)",
                    "kube-public – Publicly accessible resources",
                    "kube-node-lease – Node heartbeat data",
                ]),
                ("content", "Labels & Selectors", [
                    "Labels are key=value pairs attached to resources",
                    "They connect Services to Pods, Deployments to Pods, etc.",
                    ">> kubectl get pods --show-labels",
                    ">> kubectl get pods -l app=nginx",
                    ">> kubectl get nodes --show-labels",
                ]),
                ("table", "Essential Commands Summary", ["Action", "Command"],
                 [["View nodes", "kubectl get nodes -o wide"],
                  ["Pods in all NS", "kubectl get pods -A"],
                  ["Diagnose pod", "kubectl describe pod <name>"],
                  ["View logs", "kubectl logs <pod>"],
                  ["Shell into pod", "kubectl exec -it <pod> -- /bin/sh"],
                  ["Recent events", "kubectl get events --sort-by='.lastTimestamp'"]]),
            ],
            "lab_title": "Application Down + Scavenger Hunt",
            "lab_desc": "Fix a broken web-app (0 endpoints), then prove your kubectl skills with a scavenger hunt.",
            "lab_cmd": "lab-01.sh",
        },
        "es": {
            "title": "Fundamentos de kubectl y Arquitectura AKS",
            "subtitle": "Aprende los comandos esenciales y cómo funcionan los clusters AKS",
            "slides": [
                ("section", "Arquitectura AKS", "Entendiendo los componentes"),
                ("content", "Control Plane (Administrado por Azure)", [
                    "API Server – Punto de entrada para todas las operaciones kubectl",
                    "etcd – Base de datos que guarda el estado del cluster",
                    "Scheduler – Decide en qué nodo colocar cada pod",
                    "Controller Manager – Mantiene el estado deseado",
                    "Cloud Controller Manager – Integración con Azure",
                ]),
                ("content", "Worker Nodes (Nodos)", [
                    "kubelet – Agente que ejecuta pods en cada nodo",
                    "kube-proxy – Configura reglas de red en cada nodo",
                    "Container Runtime – Ejecuta los contenedores",
                    "CoreDNS – Resolución DNS interna del cluster",
                ]),
                ("section", "Conexión al Cluster", ""),
                ("code", "Conectarse al AKS", [
                    "# Login a Azure",
                    "az login",
                    "",
                    "# Obtener credenciales del cluster",
                    "az aks get-credentials --resource-group <rg> --name <cluster>",
                    "",
                    "# Verificar conexión",
                    "kubectl cluster-info",
                ]),
                ("section", "Comandos Esenciales de kubectl", ""),
                ("code", "Ver Recursos", [
                    "# Nodos",
                    "kubectl get nodes -o wide",
                    "",
                    "# Pods (todos los namespaces)",
                    "kubectl get pods -A",
                    "",
                    "# Todo en un namespace",
                    "kubectl get all -n <namespace>",
                ]),
                ("code", "Diagnosticar Recursos", [
                    "# Describe muestra EVENTOS - clave para troubleshooting",
                    "kubectl describe pod <nombre>",
                    "kubectl describe node <nombre>",
                    "",
                    "# Logs",
                    "kubectl logs <pod>",
                    "kubectl logs <pod> --previous",
                    "",
                    "# Ejecutar comando dentro de un pod",
                    "kubectl exec -it <pod> -- /bin/sh",
                ]),
                ("content", "Namespaces", [
                    "default – Donde van los recursos si no especificas namespace",
                    "kube-system – Componentes del sistema (CoreDNS, kube-proxy, etc.)",
                    "kube-public – Recursos públicos",
                    "kube-node-lease – Heartbeats de los nodos",
                ]),
                ("content", "Labels y Selectors", [
                    "Los labels son pares key=value asignados a recursos",
                    "Son la base de cómo Kubernetes conecta todo",
                    ">> kubectl get pods --show-labels",
                    ">> kubectl get pods -l app=nginx",
                    ">> kubectl get nodes --show-labels",
                ]),
                ("table", "Resumen de Comandos", ["Acción", "Comando"],
                 [["Ver nodos", "kubectl get nodes -o wide"],
                  ["Pods en todo el cluster", "kubectl get pods -A"],
                  ["Diagnosticar pod", "kubectl describe pod <nombre>"],
                  ["Ver logs", "kubectl logs <pod>"],
                  ["Shell en un pod", "kubectl exec -it <pod> -- /bin/sh"],
                  ["Eventos recientes", "kubectl get events --sort-by='.lastTimestamp'"]]),
            ],
            "lab_title": "Aplicación Caída + Búsqueda del Tesoro",
            "lab_desc": "Repara la web-app (0 endpoints) y demuestra tus habilidades kubectl con una búsqueda del tesoro.",
            "lab_cmd": "lab-01.sh",
        },
    },
    2: {
        "en": {
            "title": "Pods & Containers",
            "subtitle": "Understand pod lifecycle, diagnose image errors, and read events",
            "slides": [
                ("section", "What is a Pod?", "The smallest unit in Kubernetes"),
                ("content", "Pod Basics", [
                    "A Pod contains one or more containers sharing network and storage",
                    "Pods are ephemeral – when they die, data is lost",
                    "Each pod gets its own IP address within the cluster",
                    "Defined using YAML manifests with apiVersion, kind, metadata, spec",
                ]),
                ("table", "Pod Lifecycle States", ["State", "Meaning"],
                 [["Pending", "Waiting for node assignment or image pull"],
                  ["Running", "At least one container is running"],
                  ["Succeeded", "All containers exited with code 0"],
                  ["Failed", "At least one container exited with error"],
                  ["CrashLoopBackOff", "Container keeps crashing and restarting"],
                  ["ImagePullBackOff", "Cannot download the container image"]]),
                ("section", "Container Images", ""),
                ("content", "Image Format & Errors", [
                    "Format: registry/repository:tag",
                    ">> nginx:1.25 (Docker Hub implicit)",
                    ">> mcr.microsoft.com/azuredocs/aci-helloworld (MCR)",
                    ">> myacr.azurecr.io/myapp:v2 (Azure Container Registry)",
                    "",
                    "ErrImagePull – Image or tag doesn't exist, or no permissions",
                    "ImagePullBackOff – Retries failing with exponential backoff",
                ]),
                ("section", "Diagnosing Pods", "The 3-step process"),
                ("code", "Step 1: View Status", [
                    "kubectl get pods",
                    "# NAME       READY   STATUS             RESTARTS   AGE",
                    "# web-app    0/1     ImagePullBackOff   0          5m",
                ]),
                ("code", "Step 2: Describe (Always Check Events!)", [
                    "kubectl describe pod web-app",
                    "",
                    "# Events:",
                    "#  Warning  Failed  kubelet  Failed to pull image \"nginx:99.99\"",
                    "#  Warning  Failed  kubelet  Error: ErrImagePull",
                    "#  Normal   BackOff kubelet  Back-off pulling image",
                ]),
                ("code", "Step 3: Logs", [
                    "# Current logs",
                    "kubectl logs web-app",
                    "",
                    "# Previous container logs (after crash)",
                    "kubectl logs web-app --previous",
                    "",
                    "# Logs by label",
                    "kubectl logs -l app=web-app --all-containers",
                ]),
                ("content", "Resources: Requests & Limits", [
                    "requests – Minimum guaranteed resources (scheduler uses this)",
                    "limits – Maximum allowed (exceeding memory → OOMKilled)",
                    ">> kubectl top pods",
                    ">> kubectl top nodes",
                ]),
                ("table", "Diagnosis Quick Reference", ["To diagnose...", "Use..."],
                 [["Pod status", "kubectl get pods"],
                  ["Why it failed", "kubectl describe pod <name> → Events"],
                  ["App output", "kubectl logs <pod>"],
                  ["Previous crash", "kubectl logs <pod> --previous"],
                  ["Test from inside", "kubectl exec -it <pod> -- /bin/sh"]]),
            ],
            "lab_title": "Fix ImagePullBackOff",
            "lab_desc": "A pod has a wrong image tag. Diagnose and fix the issue.",
            "lab_cmd": "lab-02.sh",
        },
        "es": {
            "title": "Pods y Contenedores",
            "subtitle": "Ciclo de vida de pods, diagnosticar errores de imagen y leer eventos",
            "slides": [
                ("section", "¿Qué es un Pod?", "La unidad más pequeña en Kubernetes"),
                ("content", "Fundamentos del Pod", [
                    "Un Pod contiene uno o más contenedores que comparten red y almacenamiento",
                    "Los pods son efímeros – cuando mueren, los datos se pierden",
                    "Cada pod recibe su propia IP dentro del cluster",
                    "Se definen con manifiestos YAML: apiVersion, kind, metadata, spec",
                ]),
                ("table", "Estados del Pod", ["Estado", "Significado"],
                 [["Pending", "Esperando asignación a nodo o descarga de imagen"],
                  ["Running", "Al menos un contenedor corriendo"],
                  ["Succeeded", "Todos los contenedores terminaron exitosamente"],
                  ["Failed", "Al menos un contenedor falló"],
                  ["CrashLoopBackOff", "El contenedor se reinicia constantemente"],
                  ["ImagePullBackOff", "No puede descargar la imagen"]]),
                ("section", "Imágenes de Contenedores", ""),
                ("content", "Formato y Errores de Imagen", [
                    "Formato: registro/repositorio:tag",
                    ">> nginx:1.25 (Docker Hub implícito)",
                    ">> mcr.microsoft.com/azuredocs/aci-helloworld (MCR)",
                    ">> myacr.azurecr.io/myapp:v2 (Azure Container Registry)",
                    "",
                    "ErrImagePull – Imagen o tag no existe, o sin permisos",
                    "ImagePullBackOff – Reintentos fallando con backoff exponencial",
                ]),
                ("section", "Diagnóstico de Pods", "El proceso de 3 pasos"),
                ("code", "Paso 1: Ver Estado", [
                    "kubectl get pods",
                    "# NAME       READY   STATUS             RESTARTS   AGE",
                    "# web-app    0/1     ImagePullBackOff   0          5m",
                ]),
                ("code", "Paso 2: Describe (¡Siempre Revisar Events!)", [
                    "kubectl describe pod web-app",
                    "",
                    "# Events:",
                    "#  Warning  Failed  kubelet  Failed to pull image \"nginx:99.99\"",
                    "#  Warning  Failed  kubelet  Error: ErrImagePull",
                    "#  Normal   BackOff kubelet  Back-off pulling image",
                ]),
                ("code", "Paso 3: Logs", [
                    "# Logs actuales",
                    "kubectl logs web-app",
                    "",
                    "# Logs del container anterior (crasheado)",
                    "kubectl logs web-app --previous",
                    "",
                    "# Logs por label",
                    "kubectl logs -l app=web-app --all-containers",
                ]),
                ("content", "Resources: Requests y Limits", [
                    "requests – Mínimo garantizado (el scheduler usa esto)",
                    "limits – Máximo permitido (exceder memoria → OOMKilled)",
                    ">> kubectl top pods",
                    ">> kubectl top nodes",
                ]),
                ("table", "Referencia Rápida de Diagnóstico", ["Para diagnosticar...", "Usa..."],
                 [["Estado del pod", "kubectl get pods"],
                  ["Por qué falló", "kubectl describe pod <nombre> → Events"],
                  ["Salida de la app", "kubectl logs <pod>"],
                  ["Crash anterior", "kubectl logs <pod> --previous"],
                  ["Probar desde dentro", "kubectl exec -it <pod> -- /bin/sh"]]),
            ],
            "lab_title": "Arreglar ImagePullBackOff",
            "lab_desc": "Un pod tiene un tag de imagen incorrecta. Diagnostica y arregla el problema.",
            "lab_cmd": "lab-02.sh",
        },
    },
    3: {
        "en": {
            "title": "Deployments & ReplicaSets",
            "subtitle": "Rolling updates, rollbacks, and health checks",
            "slides": [
                ("section", "What is a Deployment?", "The standard way to run applications"),
                ("content", "Deployment → ReplicaSet → Pods", [
                    "Deployment manages desired state declaratively",
                    "ReplicaSet ensures the right number of pod replicas",
                    "Pods are the actual running containers",
                    "Deployment creates new ReplicaSets on each update",
                ]),
                ("code", "Basic Commands", [
                    "# Scale",
                    "kubectl scale deployment web-app --replicas=5",
                    "",
                    "# Update image (triggers rolling update)",
                    "kubectl set image deployment/web-app nginx=nginx:1.26",
                    "",
                    "# Check rollout status",
                    "kubectl rollout status deployment/web-app",
                ]),
                ("section", "Rolling Updates & Rollbacks", ""),
                ("code", "Rollback Commands", [
                    "# View history",
                    "kubectl rollout history deployment/web-app",
                    "",
                    "# Undo last update",
                    "kubectl rollout undo deployment/web-app",
                    "",
                    "# Rollback to specific revision",
                    "kubectl rollout undo deployment/web-app --to-revision=2",
                ]),
                ("section", "Health Checks (Probes)", "Critical for stability"),
                ("table", "Probe Types", ["Probe", "Purpose", "On Failure"],
                 [["livenessProbe", "Is the container alive?", "kubelet restarts it"],
                  ["readinessProbe", "Can it receive traffic?", "Removed from Service"],
                  ["startupProbe", "Has it finished starting?", "Blocks other probes"]]),
                ("content", "Common Probe Mistakes", [
                    "Wrong port in probe → CrashLoopBackOff",
                    "Non-existent path → CrashLoopBackOff",
                    "initialDelaySeconds too low → App killed before starting",
                    "No readinessProbe → Traffic hits unready pods",
                ]),
                ("table", "Commands Summary", ["Action", "Command"],
                 [["View deployment", "kubectl get deploy web-app"],
                  ["View pods", "kubectl get pods -l app=web-app"],
                  ["Version history", "kubectl rollout history deploy/web-app"],
                  ["Rollback", "kubectl rollout undo deploy/web-app"],
                  ["Scale", "kubectl scale deploy/web-app --replicas=5"]]),
            ],
            "lab_title": "Failed Rollout",
            "lab_desc": "A deployment has a failed rollout due to a bad liveness probe. Fix it or rollback.",
            "lab_cmd": "lab-03.sh",
        },
        "es": {
            "title": "Deployments y ReplicaSets",
            "subtitle": "Rolling updates, rollbacks y health checks",
            "slides": [
                ("section", "¿Qué es un Deployment?", "La forma estándar de desplegar aplicaciones"),
                ("content", "Deployment → ReplicaSet → Pods", [
                    "Deployment maneja el estado deseado declarativamente",
                    "ReplicaSet asegura la cantidad correcta de réplicas",
                    "Pods son los contenedores corriendo",
                    "El Deployment crea nuevos ReplicaSets con cada actualización",
                ]),
                ("code", "Comandos Básicos", [
                    "# Escalar",
                    "kubectl scale deployment web-app --replicas=5",
                    "",
                    "# Actualizar imagen (dispara rolling update)",
                    "kubectl set image deployment/web-app nginx=nginx:1.26",
                    "",
                    "# Ver progreso del rollout",
                    "kubectl rollout status deployment/web-app",
                ]),
                ("section", "Rolling Updates y Rollbacks", ""),
                ("code", "Comandos de Rollback", [
                    "# Ver historial",
                    "kubectl rollout history deployment/web-app",
                    "",
                    "# Revertir al anterior",
                    "kubectl rollout undo deployment/web-app",
                    "",
                    "# Revertir a revisión específica",
                    "kubectl rollout undo deployment/web-app --to-revision=2",
                ]),
                ("section", "Health Checks (Probes)", "Cruciales para la estabilidad"),
                ("table", "Tipos de Probes", ["Probe", "Propósito", "Si Falla"],
                 [["livenessProbe", "¿El container está vivo?", "kubelet lo reinicia"],
                  ["readinessProbe", "¿Puede recibir tráfico?", "Se saca del Service"],
                  ["startupProbe", "¿Terminó de arrancar?", "Bloquea otros probes"]]),
                ("content", "Errores Comunes con Probes", [
                    "Puerto incorrecto en probe → CrashLoopBackOff",
                    "Path inexistente → CrashLoopBackOff",
                    "initialDelaySeconds muy bajo → App se mata antes de arrancar",
                    "Sin readinessProbe → Tráfico llega a pods no listos",
                ]),
                ("table", "Resumen de Comandos", ["Acción", "Comando"],
                 [["Ver deployment", "kubectl get deploy web-app"],
                  ["Ver pods", "kubectl get pods -l app=web-app"],
                  ["Historial", "kubectl rollout history deploy/web-app"],
                  ["Rollback", "kubectl rollout undo deploy/web-app"],
                  ["Escalar", "kubectl scale deploy/web-app --replicas=5"]]),
            ],
            "lab_title": "Rollout Fallido",
            "lab_desc": "Un deployment tiene un rollout fallido por un liveness probe mal configurado. Arréglalo o haz rollback.",
            "lab_cmd": "lab-03.sh",
        },
    },
    4: {
        "en": {
            "title": "Services & Networking",
            "subtitle": "How pods communicate and expose traffic",
            "slides": [
                ("section", "What is a Service?", "Stable endpoint for ephemeral pods"),
                ("table", "Service Types", ["Type", "Access", "Use Case"],
                 [["ClusterIP", "Internal only", "Backend APIs, databases"],
                  ["NodePort", "Port on each node (30000-32767)", "Dev/test access"],
                  ["LoadBalancer", "Azure LB with public IP", "Production external access"]]),
                ("section", "Internal DNS", "How pods find each other"),
                ("content", "DNS Resolution", [
                    "Format: <service>.<namespace>.svc.cluster.local",
                    ">> Same namespace: curl http://backend-svc",
                    ">> Other namespace: curl http://backend-svc.production",
                    "CoreDNS runs in kube-system and handles all DNS",
                ]),
                ("code", "Testing DNS", [
                    "# Temporary debug pod",
                    "kubectl run debug --image=busybox:1.36 --rm -it -- /bin/sh",
                    "",
                    "# Inside the pod:",
                    "nslookup backend-svc",
                    "wget -qO- http://backend-svc --timeout=5",
                ]),
                ("section", "Troubleshooting Services", ""),
                ("content", "Service Has No Endpoints", [
                    "Endpoints are the pod IPs that a Service routes traffic to",
                    "Empty endpoints = Service can't find matching pods",
                    ">> kubectl get endpoints my-svc",
                    "",
                    "Check: Does the Service selector match the pod labels?",
                    ">> kubectl describe svc my-svc | grep Selector",
                    ">> kubectl get pods --show-labels",
                ]),
                ("table", "Common Issues", ["Symptom", "Likely Cause"],
                 [["Service no response", "Endpoints empty (selector mismatch)"],
                  ["Connection refused", "Wrong targetPort vs container port"],
                  ["DNS not resolving", "CoreDNS down or NetworkPolicy blocking"],
                  ["LoadBalancer no IP", "Quota exhausted or permissions issue"]]),
                ("code", "Connectivity Checklist", [
                    "# 1. Service exists with IP?",
                    "kubectl get svc web-svc",
                    "# 2. Has endpoints?",
                    "kubectl get endpoints web-svc",
                    "# 3. Pods running?",
                    "kubectl get pods -l app=web",
                    "# 4. Selector matches labels?",
                    "kubectl describe svc web-svc | grep Selector",
                    "# 5. Can connect from another pod?",
                    "kubectl exec debug -- wget -qO- http://web-svc --timeout=5",
                ]),
            ],
            "lab_title": "Service Disconnect",
            "lab_desc": "A Service can't connect to its pods. Find the mismatch and fix it.",
            "lab_cmd": "lab-04.sh",
        },
        "es": {
            "title": "Services y Networking",
            "subtitle": "Cómo se comunican los pods y exponen tráfico",
            "slides": [
                ("section", "¿Qué es un Service?", "Endpoint estable para pods efímeros"),
                ("table", "Tipos de Service", ["Tipo", "Acceso", "Caso de Uso"],
                 [["ClusterIP", "Solo interno", "APIs backend, bases de datos"],
                  ["NodePort", "Puerto en cada nodo (30000-32767)", "Dev/test"],
                  ["LoadBalancer", "Azure LB con IP pública", "Acceso externo producción"]]),
                ("section", "DNS Interno", "Cómo los pods se encuentran"),
                ("content", "Resolución DNS", [
                    "Formato: <servicio>.<namespace>.svc.cluster.local",
                    ">> Mismo namespace: curl http://backend-svc",
                    ">> Otro namespace: curl http://backend-svc.production",
                    "CoreDNS corre en kube-system y maneja todo el DNS",
                ]),
                ("code", "Probar DNS", [
                    "# Pod temporal de debug",
                    "kubectl run debug --image=busybox:1.36 --rm -it -- /bin/sh",
                    "",
                    "# Dentro del pod:",
                    "nslookup backend-svc",
                    "wget -qO- http://backend-svc --timeout=5",
                ]),
                ("section", "Troubleshooting de Services", ""),
                ("content", "Service Sin Endpoints", [
                    "Endpoints son las IPs de pods a las que el Service envía tráfico",
                    "Endpoints vacío = Service no encuentra pods",
                    ">> kubectl get endpoints my-svc",
                    "",
                    "Verificar: ¿El selector del Service coincide con labels de los pods?",
                    ">> kubectl describe svc my-svc | grep Selector",
                    ">> kubectl get pods --show-labels",
                ]),
                ("table", "Problemas Comunes", ["Síntoma", "Causa Probable"],
                 [["Service sin respuesta", "Endpoints vacío (selector no coincide)"],
                  ["Connection refused", "targetPort incorrecto vs puerto del container"],
                  ["DNS no resuelve", "CoreDNS caído o NetworkPolicy bloqueando"],
                  ["LoadBalancer sin IP", "Cuota agotada o problema de permisos"]]),
                ("code", "Checklist de Conectividad", [
                    "# 1. ¿Service existe con IP?",
                    "kubectl get svc web-svc",
                    "# 2. ¿Tiene endpoints?",
                    "kubectl get endpoints web-svc",
                    "# 3. ¿Pods corriendo?",
                    "kubectl get pods -l app=web",
                    "# 4. ¿Selector coincide con labels?",
                    "kubectl describe svc web-svc | grep Selector",
                    "# 5. ¿Se puede conectar desde otro pod?",
                    "kubectl exec debug -- wget -qO- http://web-svc --timeout=5",
                ]),
            ],
            "lab_title": "Service Desconectado",
            "lab_desc": "Un Service no puede conectar con sus pods. Encuentra el problema y arréglalo.",
            "lab_cmd": "lab-04.sh",
        },
    },
    5: {
        "en": {
            "title": "ConfigMaps & Secrets",
            "subtitle": "Externalize configuration and manage sensitive data",
            "slides": [
                ("section", "Why ConfigMaps & Secrets?", "Never hardcode config in images"),
                ("table", "When to Use", ["Data Type", "Resource"],
                 [["URLs, feature flags, config", "ConfigMap"],
                  ["Passwords, tokens, certs", "Secret"]]),
                ("code", "Creating ConfigMaps", [
                    "# From literal",
                    "kubectl create configmap app-config \\",
                    "  --from-literal=DATABASE_HOST=db.example.com \\",
                    "  --from-literal=LOG_LEVEL=info",
                    "",
                    "# From file",
                    "kubectl create configmap nginx-config --from-file=nginx.conf",
                    "",
                    "# View contents",
                    "kubectl describe configmap app-config",
                ]),
                ("content", "Using in Pods", [
                    "As environment variables: env.valueFrom.configMapKeyRef",
                    "As mounted files: volumes + volumeMounts",
                    "envFrom loads ALL keys as env vars at once",
                ]),
                ("code", "Creating Secrets", [
                    "kubectl create secret generic db-creds \\",
                    "  --from-literal=username=admin \\",
                    "  --from-literal=password=S3cureP@ss!",
                    "",
                    "# Values stored in base64 (NOT encrypted!)",
                    "kubectl get secret db-creds -o yaml",
                    "",
                    "# Decode",
                    "echo 'YWRtaW4=' | base64 --decode  # → admin",
                ]),
                ("table", "Common Errors", ["Error", "Cause"],
                 [["CreateContainerConfigError", "ConfigMap or Secret doesn't exist"],
                  ["Pod won't start", "Referenced key doesn't exist in ConfigMap"],
                  ["Empty env var", "Wrong key name or optional: true"]]),
                ("code", "Diagnosis", [
                    "# Check if ConfigMap exists",
                    "kubectl get configmap app-config",
                    "",
                    "# Check keys",
                    "kubectl describe configmap app-config",
                    "",
                    "# Check env vars inside running pod",
                    "kubectl exec <pod> -- env | sort",
                ]),
            ],
            "lab_title": "Missing Configuration",
            "lab_desc": "An app crashes because of missing ConfigMap/Secret. Create the missing resources.",
            "lab_cmd": "lab-05.sh",
        },
        "es": {
            "title": "ConfigMaps y Secrets",
            "subtitle": "Externalizar configuración y manejar datos sensibles",
            "slides": [
                ("section", "¿Por qué ConfigMaps y Secrets?", "Nunca hardcodear config en las imágenes"),
                ("table", "Cuándo Usar", ["Tipo de Dato", "Recurso"],
                 [["URLs, feature flags, config", "ConfigMap"],
                  ["Contraseñas, tokens, certificados", "Secret"]]),
                ("code", "Crear ConfigMaps", [
                    "# Desde literal",
                    "kubectl create configmap app-config \\",
                    "  --from-literal=DATABASE_HOST=db.example.com \\",
                    "  --from-literal=LOG_LEVEL=info",
                    "",
                    "# Desde archivo",
                    "kubectl create configmap nginx-config --from-file=nginx.conf",
                    "",
                    "# Ver contenido",
                    "kubectl describe configmap app-config",
                ]),
                ("content", "Usar en Pods", [
                    "Como variables de entorno: env.valueFrom.configMapKeyRef",
                    "Como archivos montados: volumes + volumeMounts",
                    "envFrom carga TODAS las keys como env vars de una vez",
                ]),
                ("code", "Crear Secrets", [
                    "kubectl create secret generic db-creds \\",
                    "  --from-literal=username=admin \\",
                    "  --from-literal=password=S3cureP@ss!",
                    "",
                    "# Valores en base64 (NO es encriptación!)",
                    "kubectl get secret db-creds -o yaml",
                    "",
                    "# Decodificar",
                    "echo 'YWRtaW4=' | base64 --decode  # → admin",
                ]),
                ("table", "Errores Comunes", ["Error", "Causa"],
                 [["CreateContainerConfigError", "ConfigMap o Secret no existe"],
                  ["Pod no arranca", "Key referenciada no existe"],
                  ["Env var vacía", "Nombre de key incorrecto u optional: true"]]),
                ("code", "Diagnóstico", [
                    "# Verificar que ConfigMap existe",
                    "kubectl get configmap app-config",
                    "",
                    "# Ver las keys",
                    "kubectl describe configmap app-config",
                    "",
                    "# Ver env vars dentro del pod",
                    "kubectl exec <pod> -- env | sort",
                ]),
            ],
            "lab_title": "Configuración Faltante",
            "lab_desc": "Una app falla porque le falta ConfigMap/Secret. Crea los recursos faltantes.",
            "lab_cmd": "lab-05.sh",
        },
    },
    6: {
        "en": {
            "title": "Storage & Volumes",
            "subtitle": "Persistent data with PV, PVC, and StorageClasses in AKS",
            "slides": [
                ("section", "The Problem", "Containers are ephemeral – data is lost on restart"),
                ("content", "Key Concepts", [
                    "StorageClass – Defines HOW storage is provisioned",
                    "PersistentVolume (PV) – A real provisioned disk",
                    "PersistentVolumeClaim (PVC) – A pod's request for storage",
                    "Flow: Pod → PVC → PV → Azure Disk / Azure Files",
                ]),
                ("table", "StorageClasses in AKS", ["StorageClass", "Type", "Access"],
                 [["managed-csi", "Azure Disk (SSD)", "ReadWriteOnce (RWO)"],
                  ["managed-csi-premium", "Azure Premium Disk", "ReadWriteOnce (RWO)"],
                  ["azurefile-csi", "Azure Files", "ReadWriteMany (RWX)"],
                  ["azurefile-csi-premium", "Azure Files Premium", "ReadWriteMany (RWX)"]]),
                ("content", "Access Modes", [
                    "ReadWriteOnce (RWO) – One node can mount read/write",
                    "ReadWriteMany (RWX) – Multiple nodes can mount read/write",
                    "Azure Disks only support RWO. Need RWX → use Azure Files",
                ]),
                ("table", "PVC Stuck in Pending", ["Cause", "Diagnosis"],
                 [["StorageClass doesn't exist", "kubectl get sc → check name"],
                  ["Zone mismatch", "Disk in zone 1, node in zone 2"],
                  ["Quota exceeded", "kubectl describe pvc → Events"],
                  ["Wrong accessModes", "RWX with Azure Disk (only supports RWO)"],
                  ["WaitForFirstConsumer", "Normal – binds when a pod uses it"]]),
                ("table", "Azure Disk vs Azure Files", ["Feature", "Azure Disk", "Azure Files"],
                 [["Access", "RWO (one node)", "RWX (multi node)"],
                  ["Performance", "High (SSD/Premium)", "Medium"],
                  ["Use case", "Databases", "Shared files"],
                  ["StorageClass", "managed-csi", "azurefile-csi"]]),
                ("code", "Diagnosis Commands", [
                    "kubectl get pvc                    # Status",
                    "kubectl describe pvc <name>        # Events!",
                    "kubectl get pv                     # Provisioned disks",
                    "kubectl get sc                     # StorageClasses",
                ]),
            ],
            "lab_title": "PVC Stuck in Pending",
            "lab_desc": "A PVC is stuck in Pending state. Find the cause and fix it.",
            "lab_cmd": "lab-06.sh",
        },
        "es": {
            "title": "Storage y Volúmenes",
            "subtitle": "Datos persistentes con PV, PVC y StorageClasses en AKS",
            "slides": [
                ("section", "El Problema", "Los containers son efímeros – datos se pierden al reiniciar"),
                ("content", "Conceptos Clave", [
                    "StorageClass – Define CÓMO se crea el almacenamiento",
                    "PersistentVolume (PV) – Un disco real provisionado",
                    "PersistentVolumeClaim (PVC) – La solicitud de almacenamiento del pod",
                    "Flujo: Pod → PVC → PV → Azure Disk / Azure Files",
                ]),
                ("table", "StorageClasses en AKS", ["StorageClass", "Tipo", "Acceso"],
                 [["managed-csi", "Azure Disk (SSD)", "ReadWriteOnce (RWO)"],
                  ["managed-csi-premium", "Azure Premium Disk", "ReadWriteOnce (RWO)"],
                  ["azurefile-csi", "Azure Files", "ReadWriteMany (RWX)"],
                  ["azurefile-csi-premium", "Azure Files Premium", "ReadWriteMany (RWX)"]]),
                ("content", "Modos de Acceso", [
                    "ReadWriteOnce (RWO) – Un solo nodo puede montar lectura/escritura",
                    "ReadWriteMany (RWX) – Múltiples nodos pueden montar",
                    "Azure Disks solo soportan RWO. ¿Necesitas RWX? → Azure Files",
                ]),
                ("table", "PVC Atascado en Pending", ["Causa", "Diagnóstico"],
                 [["StorageClass no existe", "kubectl get sc → verificar nombre"],
                  ["Zona incorrecta", "Disco en zona 1, nodo en zona 2"],
                  ["Cuota excedida", "kubectl describe pvc → Events"],
                  ["accessModes incompatible", "RWX con Azure Disk (solo soporta RWO)"],
                  ["WaitForFirstConsumer", "Normal – se vincula cuando un pod lo usa"]]),
                ("table", "Azure Disk vs Azure Files", ["Característica", "Azure Disk", "Azure Files"],
                 [["Acceso", "RWO (un nodo)", "RWX (multi nodo)"],
                  ["Performance", "Alta (SSD/Premium)", "Media"],
                  ["Caso de uso", "Bases de datos", "Archivos compartidos"],
                  ["StorageClass", "managed-csi", "azurefile-csi"]]),
                ("code", "Comandos de Diagnóstico", [
                    "kubectl get pvc                    # Estado",
                    "kubectl describe pvc <nombre>      # ¡Events!",
                    "kubectl get pv                     # Discos provisionados",
                    "kubectl get sc                     # StorageClasses",
                ]),
            ],
            "lab_title": "PVC Atascado en Pending",
            "lab_desc": "Un PVC está atascado en Pending. Encuentra la causa y arréglalo.",
            "lab_cmd": "lab-06.sh",
        },
    },
    7: {
        "en": {
            "title": "Network Policies",
            "subtitle": "Firewalls at the pod level – controlling traffic between pods",
            "slides": [
                ("section", "What are Network Policies?", "By default all pods can talk to all pods"),
                ("content", "Key Rules", [
                    "No NetworkPolicy → everything is allowed",
                    "If a policy exists for a pod → everything NOT allowed is BLOCKED",
                    "policyTypes: [Ingress] with no rules → blocks all inbound",
                    "policyTypes: [Egress] with no rules → blocks all outbound (incl. DNS!)",
                ]),
                ("content", "Anatomy of a NetworkPolicy", [
                    "podSelector – Which pods does this policy apply to?",
                    "policyTypes – Ingress, Egress, or both",
                    "ingress.from – Who can send traffic IN",
                    "egress.to – Where can traffic go OUT",
                    "ports – Which ports are allowed",
                ]),
                ("content", "Selectors", [
                    "podSelector – Pods in the same namespace",
                    "namespaceSelector – Pods from other namespaces",
                    "ipBlock – External IP ranges (CIDR)",
                ]),
                ("section", "Critical: DNS & Egress", "Most common mistake"),
                ("content", "Always Allow DNS in Egress Policies", [
                    "DNS uses UDP/TCP port 53 to CoreDNS in kube-system",
                    "Blocking egress without DNS exception = nothing resolves",
                    "Add explicit egress rule for port 53 to any namespace",
                ]),
                ("code", "Diagnosis", [
                    "# List policies",
                    "kubectl get networkpolicy",
                    "",
                    "# Detail",
                    "kubectl describe netpol <name>",
                    "",
                    "# Test connectivity",
                    "kubectl exec <pod-a> -- wget -qO- http://<svc> --timeout=5",
                    "",
                    "# Test DNS",
                    "kubectl exec <pod> -- nslookup kubernetes.default",
                    "",
                    "# Check pod labels (match selectors?)",
                    "kubectl get pods --show-labels",
                ]),
                ("table", "Checklist", ["Check", "Command"],
                 [["Any policies exist?", "kubectl get netpol"],
                  ["Selector matches pod?", "Compare labels vs podSelector"],
                  ["DNS allowed in egress?", "Check for port 53 rule"],
                  ["From labels correct?", "kubectl get pods --show-labels"]]),
            ],
            "lab_title": "Blocked Traffic",
            "lab_desc": "A NetworkPolicy is blocking legitimate traffic. Find and fix the misconfigured rule.",
            "lab_cmd": "lab-07.sh",
        },
        "es": {
            "title": "Network Policies",
            "subtitle": "Firewalls a nivel de pod – controlar tráfico entre pods",
            "slides": [
                ("section", "¿Qué son las Network Policies?", "Por defecto todos los pods pueden hablar con todos"),
                ("content", "Reglas Clave", [
                    "Sin NetworkPolicy → todo está permitido",
                    "Si existe una policy para un pod → todo lo NO permitido se BLOQUEA",
                    "policyTypes: [Ingress] sin reglas → bloquea todo el ingreso",
                    "policyTypes: [Egress] sin reglas → bloquea todo el egreso (¡incluyendo DNS!)",
                ]),
                ("content", "Anatomía de una NetworkPolicy", [
                    "podSelector – ¿A qué pods aplica esta regla?",
                    "policyTypes – Ingress, Egress, o ambos",
                    "ingress.from – Quién puede enviar tráfico HACIA el pod",
                    "egress.to – A dónde puede ir el tráfico DESDE el pod",
                    "ports – Qué puertos se permiten",
                ]),
                ("content", "Selectores", [
                    "podSelector – Pods en el mismo namespace",
                    "namespaceSelector – Pods de otros namespaces",
                    "ipBlock – Rangos de IPs externas (CIDR)",
                ]),
                ("section", "Crítico: DNS y Egress", "El error más común"),
                ("content", "Siempre Permitir DNS en Políticas de Egress", [
                    "DNS usa UDP/TCP puerto 53 hacia CoreDNS en kube-system",
                    "Bloquear egress sin excepción DNS = nada resuelve",
                    "Agregar regla egress explícita para puerto 53",
                ]),
                ("code", "Diagnóstico", [
                    "# Listar policies",
                    "kubectl get networkpolicy",
                    "",
                    "# Detalle",
                    "kubectl describe netpol <nombre>",
                    "",
                    "# Probar conectividad",
                    "kubectl exec <pod-a> -- wget -qO- http://<svc> --timeout=5",
                    "",
                    "# Probar DNS",
                    "kubectl exec <pod> -- nslookup kubernetes.default",
                    "",
                    "# Ver labels de pods (¿coinciden con selectores?)",
                    "kubectl get pods --show-labels",
                ]),
                ("table", "Checklist", ["Verificar", "Comando"],
                 [["¿Existen policies?", "kubectl get netpol"],
                  ["¿Selector coincide?", "Comparar labels vs podSelector"],
                  ["¿DNS permitido en egress?", "Buscar regla puerto 53"],
                  ["¿Labels from correctos?", "kubectl get pods --show-labels"]]),
            ],
            "lab_title": "Tráfico Bloqueado",
            "lab_desc": "Una NetworkPolicy está bloqueando tráfico legítimo. Encuentra y arregla la regla.",
            "lab_cmd": "lab-07.sh",
        },
    },
    8: {
        "en": {
            "title": "Node Management: Taints, Tolerations & Scheduling",
            "subtitle": "How the scheduler assigns pods to nodes",
            "slides": [
                ("section", "How Does the Scheduler Decide?", ""),
                ("content", "Scheduler Evaluation", [
                    "1. Does the node have enough resources? (CPU, memory)",
                    "2. Does the pod have a nodeSelector the node doesn't match?",
                    "3. Does the node have taints the pod doesn't tolerate?",
                    "4. Are there affinity/anti-affinity rules?",
                    "",
                    "If no node passes all filters → pod stays Pending",
                ]),
                ("section", "Taints & Tolerations", ""),
                ("content", "Taints = Repellent on Nodes", [
                    "NoSchedule – Don't schedule new pods",
                    "PreferNoSchedule – Try to avoid, not strict",
                    "NoExecute – Don't schedule + evict existing pods",
                    "",
                    "Only pods with the matching toleration can be scheduled",
                ]),
                ("code", "Taint Commands", [
                    "# Add taint",
                    "kubectl taint nodes node1 dedicated=gpu:NoSchedule",
                    "",
                    "# View taints",
                    "kubectl describe node <name> | grep -A5 Taints",
                    "",
                    "# Remove taint (trailing dash '-')",
                    "kubectl taint nodes node1 dedicated=gpu:NoSchedule-",
                    "kubectl taint nodes --all dedicated=gpu:NoSchedule-",
                ]),
                ("content", "nodeSelector", [
                    "Simplest way to choose a node",
                    "Pod spec: nodeSelector.disktype: ssd",
                    "No matching label on any node → pod stays Pending",
                    ">> kubectl get nodes --show-labels",
                    ">> kubectl label nodes node1 disktype=ssd",
                ]),
                ("code", "Cordon & Drain", [
                    "# Cordon: mark node as unschedulable",
                    "kubectl cordon node1",
                    "",
                    "# Uncordon: allow scheduling again",
                    "kubectl uncordon node1",
                    "",
                    "# Drain: move all pods to other nodes (maintenance)",
                    "kubectl drain node1 --ignore-daemonsets --delete-emptydir-data",
                ]),
                ("table", "Diagnosis", ["Events Message", "Cause", "Fix"],
                 [["didn't tolerate taint", "Node has taint", "Add toleration or remove taint"],
                  ["didn't match node selector", "Missing label", "Add label or remove selector"],
                  ["Insufficient cpu/memory", "No resources", "Scale pool or reduce requests"]]),
            ],
            "lab_title": "Node Taint Issue",
            "lab_desc": "Maintenance taints prevent pod scheduling. Find and resolve the issue.",
            "lab_cmd": "lab-08.sh",
        },
        "es": {
            "title": "Gestión de Nodos: Taints, Tolerations y Scheduling",
            "subtitle": "Cómo el scheduler asigna pods a nodos",
            "slides": [
                ("section", "¿Cómo Decide el Scheduler?", ""),
                ("content", "Evaluación del Scheduler", [
                    "1. ¿El nodo tiene recursos suficientes? (CPU, memoria)",
                    "2. ¿El pod tiene nodeSelector que el nodo no cumple?",
                    "3. ¿El nodo tiene taints que el pod no tolera?",
                    "4. ¿Hay reglas de affinity/anti-affinity?",
                    "",
                    "Si ningún nodo pasa todos los filtros → pod queda Pending",
                ]),
                ("section", "Taints y Tolerations", ""),
                ("content", "Taints = Repelente en Nodos", [
                    "NoSchedule – No agendar nuevos pods",
                    "PreferNoSchedule – Intenta evitar, no es estricto",
                    "NoExecute – No agendar + desalojar pods existentes",
                    "",
                    "Solo pods con la toleration correcta pueden ser schedulados",
                ]),
                ("code", "Comandos de Taints", [
                    "# Agregar taint",
                    "kubectl taint nodes node1 dedicated=gpu:NoSchedule",
                    "",
                    "# Ver taints",
                    "kubectl describe node <nombre> | grep -A5 Taints",
                    "",
                    "# Quitar taint (guión '-' al final)",
                    "kubectl taint nodes node1 dedicated=gpu:NoSchedule-",
                    "kubectl taint nodes --all dedicated=gpu:NoSchedule-",
                ]),
                ("content", "nodeSelector", [
                    "La forma más simple de elegir un nodo",
                    "En el pod spec: nodeSelector.disktype: ssd",
                    "Sin label en ningún nodo → pod queda Pending",
                    ">> kubectl get nodes --show-labels",
                    ">> kubectl label nodes node1 disktype=ssd",
                ]),
                ("code", "Cordon y Drain", [
                    "# Cordon: marcar nodo como no-schedulable",
                    "kubectl cordon node1",
                    "",
                    "# Uncordon: permitir scheduling de nuevo",
                    "kubectl uncordon node1",
                    "",
                    "# Drain: mover todos los pods a otros nodos",
                    "kubectl drain node1 --ignore-daemonsets --delete-emptydir-data",
                ]),
                ("table", "Diagnóstico", ["Mensaje en Events", "Causa", "Fix"],
                 [["didn't tolerate taint", "Nodo tiene taint", "Añadir toleration o quitar taint"],
                  ["didn't match node selector", "Falta label", "Añadir label o quitar selector"],
                  ["Insufficient cpu/memory", "Sin recursos", "Escalar pool o reducir requests"]]),
            ],
            "lab_title": "Problema de Taints en Nodos",
            "lab_desc": "Taints de mantenimiento impiden el scheduling. Encuentra y resuelve el problema.",
            "lab_cmd": "lab-08.sh",
        },
    },
    9: {
        "en": {
            "title": "Azure Integration: NSG, Load Balancer & Networking",
            "subtitle": "How AKS interacts with Azure infrastructure",
            "slides": [
                ("section", "AKS Azure Architecture", "The MC_ resource group"),
                ("content", "MC_ Resource Group Contents", [
                    "VMSS (Virtual Machine Scale Sets) – the nodes",
                    "VNet + Subnets – cluster network",
                    "NSG (Network Security Group) – firewall rules",
                    "Load Balancer – for LoadBalancer Services",
                    "Public IPs – external addresses",
                    "Route Table – network routing",
                ]),
                ("code", "Find the MC_ Resource Group", [
                    "# Get the managed resource group name",
                    "az aks show -g <rg> -n <cluster> \\",
                    "  --query nodeResourceGroup -o tsv",
                    "",
                    "# List all resources in it",
                    "az resource list -g MC_<rg>_<cluster>_<region> -o table",
                ]),
                ("section", "Network Security Groups (NSG)", "Azure-level firewalls"),
                ("content", "NSG Priority Rules", [
                    "Rules evaluated by priority (lowest number = highest priority)",
                    "A Deny at priority 100 overrides Allow at priority 200",
                    "AKS creates rules automatically for LB Services",
                    "Custom rules can accidentally block AKS traffic",
                ]),
                ("code", "NSG Commands", [
                    "# List NSGs",
                    "az network nsg list -g <mc-rg> -o table",
                    "",
                    "# List rules",
                    "az network nsg rule list --nsg-name <nsg> -g <mc-rg> -o table",
                    "",
                    "# Find deny rules",
                    "az network nsg rule list --nsg-name <nsg> -g <mc-rg> \\",
                    "  --query \"[?access=='Deny']\" -o table",
                    "",
                    "# Delete a rule",
                    "az network nsg rule delete --nsg-name <nsg> -g <mc-rg> -n <rule>",
                ]),
                ("table", "Common Issues", ["Symptom", "Cause"],
                 [["LB Service timeout", "NSG deny rule blocking port"],
                  ["Service no External-IP", "Public IP quota exceeded"],
                  ["Intermittent failures", "LB health probe failing"],
                  ["Can't pull images", "NSG blocking outbound 443"]]),
                ("code", "Full Diagnosis Flow", [
                    "MC_RG=$(az aks show -g <rg> -n <c> --query nodeResourceGroup -o tsv)",
                    "",
                    "# NSG rules",
                    "az network nsg list -g $MC_RG -o table",
                    "az network nsg rule list --nsg-name <nsg> -g $MC_RG -o table",
                    "",
                    "# Load Balancer",
                    "az network lb list -g $MC_RG -o table",
                    "",
                    "# Public IPs",
                    "az network public-ip list -g $MC_RG -o table",
                ]),
            ],
            "lab_title": "NSG Blocking Traffic",
            "lab_desc": "A LoadBalancer Service has an IP but HTTP times out. Find and remove the NSG deny rule.",
            "lab_cmd": "lab-09.sh",
        },
        "es": {
            "title": "Integración Azure: NSG, Load Balancer y Networking",
            "subtitle": "Cómo AKS interactúa con la infraestructura Azure",
            "slides": [
                ("section", "Arquitectura Azure de AKS", "El resource group MC_"),
                ("content", "Contenido del MC_ Resource Group", [
                    "VMSS (Virtual Machine Scale Sets) – los nodos",
                    "VNet + Subnets – red del cluster",
                    "NSG (Network Security Group) – reglas de firewall",
                    "Load Balancer – para Services tipo LoadBalancer",
                    "Public IPs – direcciones externas",
                    "Route Table – enrutamiento de red",
                ]),
                ("code", "Encontrar el MC_ Resource Group", [
                    "# Obtener el nombre del resource group managed",
                    "az aks show -g <rg> -n <cluster> \\",
                    "  --query nodeResourceGroup -o tsv",
                    "",
                    "# Listar todos los recursos",
                    "az resource list -g MC_<rg>_<cluster>_<region> -o table",
                ]),
                ("section", "Network Security Groups (NSG)", "Firewalls a nivel de Azure"),
                ("content", "Reglas de Prioridad NSG", [
                    "Las reglas se evalúan por prioridad (número más bajo = más prioritario)",
                    "Un Deny en prioridad 100 supera un Allow en prioridad 200",
                    "AKS crea reglas automáticamente para Services LB",
                    "Reglas personalizadas pueden bloquear tráfico de AKS accidentalmente",
                ]),
                ("code", "Comandos NSG", [
                    "# Listar NSGs",
                    "az network nsg list -g <mc-rg> -o table",
                    "",
                    "# Listar reglas",
                    "az network nsg rule list --nsg-name <nsg> -g <mc-rg> -o table",
                    "",
                    "# Buscar reglas deny",
                    "az network nsg rule list --nsg-name <nsg> -g <mc-rg> \\",
                    "  --query \"[?access=='Deny']\" -o table",
                    "",
                    "# Eliminar una regla",
                    "az network nsg rule delete --nsg-name <nsg> -g <mc-rg> -n <regla>",
                ]),
                ("table", "Problemas Comunes", ["Síntoma", "Causa"],
                 [["LB Service timeout", "Regla NSG deny bloqueando puerto"],
                  ["Service sin External-IP", "Cuota de IPs públicas agotada"],
                  ["Fallas intermitentes", "Health probe del LB fallando"],
                  ["No puede bajar imágenes", "NSG bloqueando outbound 443"]]),
                ("code", "Flujo de Diagnóstico Completo", [
                    "MC_RG=$(az aks show -g <rg> -n <c> --query nodeResourceGroup -o tsv)",
                    "",
                    "# Reglas NSG",
                    "az network nsg list -g $MC_RG -o table",
                    "az network nsg rule list --nsg-name <nsg> -g $MC_RG -o table",
                    "",
                    "# Load Balancer",
                    "az network lb list -g $MC_RG -o table",
                    "",
                    "# IPs públicas",
                    "az network public-ip list -g $MC_RG -o table",
                ]),
            ],
            "lab_title": "NSG Bloqueando Tráfico",
            "lab_desc": "Un Service LoadBalancer tiene IP pero HTTP da timeout. Encuentra y elimina la regla NSG deny.",
            "lab_cmd": "lab-09.sh",
        },
    },
    10: {
        "en": {
            "title": "Advanced Troubleshooting: Methodology & Complex Scenarios",
            "subtitle": "Putting it all together with the DISCOVER framework",
            "slides": [
                ("section", "The DISCOVER Framework", "A structured approach to troubleshooting"),
                ("content", "DISCOVER Steps", [
                    "D – Define: What is the exact symptom?",
                    "I – Investigate: Gather info (logs, events, describe)",
                    "S – Scope: One pod? One node? Whole cluster?",
                    "C – Compare: What changed? When did it start?",
                    "O – Options: List possible causes",
                    "V – Verify: Test your hypothesis",
                    "E – Execute: Apply the fix",
                    "R – Review: Confirm it's resolved",
                ]),
                ("section", "Advanced Commands", ""),
                ("code", "Events & Logs", [
                    "# All warnings across cluster",
                    "kubectl get events --field-selector type=Warning -A",
                    "",
                    "# Logs from all pods with a label",
                    "kubectl logs -l app=web-app --all-containers",
                    "",
                    "# Logs from last 5 minutes",
                    "kubectl logs <pod> --since=5m",
                    "",
                    "# Previous crashed container",
                    "kubectl logs <pod> --previous",
                ]),
                ("code", "Advanced JSON Queries", [
                    "# Pods NOT Running",
                    "kubectl get pods -A -o json | jq '.items[] |",
                    "  select(.status.phase != \"Running\") |",
                    "  {name: .metadata.name, status: .status.phase}'",
                    "",
                    "# Pods with restart count > 5",
                    "kubectl get pods -A -o json | jq '.items[] |",
                    "  select(.status.containerStatuses[]?.restartCount > 5)'",
                ]),
                ("section", "Complex Scenarios", ""),
                ("content", "Intermittent App Failures", [
                    "1. Pods restarting? → check restartCount",
                    "2. Resource limits hit? → kubectl top pods",
                    "3. OOMKilled? → describe pod → State",
                    "4. Readiness probe flapping? → describe pod → Conditions",
                ]),
                ("content", "Node NotReady", [
                    "1. Node conditions: kubectl describe node → Conditions",
                    "2. Kubelet status: systemctl status kubelet",
                    "3. Disk full? → Allocated resources section",
                    "4. Node events: kubectl get events for node",
                ]),
                ("code", "Quick Troubleshooting Checklist", [
                    "kubectl get pods -A              # All pods status",
                    "kubectl get nodes -o wide        # All nodes status",
                    "kubectl get events -A --sort-by='.lastTimestamp'",
                    "kubectl describe pod <name>      # Events!",
                    "kubectl logs <pod> --previous    # Crash logs",
                    "kubectl get svc,endpoints        # Service connected?",
                    "kubectl get netpol               # Policies blocking?",
                    "az aks show --query provisioningState  # Cluster OK?",
                ]),
                ("table", "Course Summary", ["Lesson", "Key Tool"],
                 [["01 kubectl", "get, describe, logs"],
                  ["02 Pods", "describe pod → Events"],
                  ["03 Deployments", "rollout status, rollout undo"],
                  ["04 Services", "get endpoints, labels matching"],
                  ["05 ConfigMaps", "describe configmap/secret"],
                  ["06 Storage", "describe pvc, StorageClass"],
                  ["07 NetPol", "get netpol, selector matching"],
                  ["08 Nodes", "taints, labels, describe node"],
                  ["09 Azure", "az network nsg rule list"],
                  ["10 Advanced", "DISCOVER methodology"]]),
            ],
            "lab_title": "Multi-Problem Challenge",
            "lab_desc": "A cluster has multiple issues combined. Apply everything you learned to fix them all.",
            "lab_cmd": "lab-10.sh",
        },
        "es": {
            "title": "Troubleshooting Avanzado: Metodología y Escenarios Complejos",
            "subtitle": "Todo junto con el framework DISCOVER",
            "slides": [
                ("section", "El Framework DISCOVER", "Enfoque estructurado de troubleshooting"),
                ("content", "Pasos DISCOVER", [
                    "D – Define: ¿Cuál es el síntoma exacto?",
                    "I – Investigate: Recopilar info (logs, events, describe)",
                    "S – Scope: ¿Un pod? ¿Un nodo? ¿Todo el cluster?",
                    "C – Compare: ¿Qué cambió? ¿Cuándo empezó?",
                    "O – Options: Listar posibles causas",
                    "V – Verify: Probar la hipótesis",
                    "E – Execute: Aplicar el fix",
                    "R – Review: Confirmar que se resolvió",
                ]),
                ("section", "Comandos Avanzados", ""),
                ("code", "Eventos y Logs", [
                    "# Todos los warnings del cluster",
                    "kubectl get events --field-selector type=Warning -A",
                    "",
                    "# Logs de todos los pods con un label",
                    "kubectl logs -l app=web-app --all-containers",
                    "",
                    "# Logs de los últimos 5 minutos",
                    "kubectl logs <pod> --since=5m",
                    "",
                    "# Container anterior crasheado",
                    "kubectl logs <pod> --previous",
                ]),
                ("code", "Queries JSON Avanzados", [
                    "# Pods que NO están Running",
                    "kubectl get pods -A -o json | jq '.items[] |",
                    "  select(.status.phase != \"Running\") |",
                    "  {name: .metadata.name, status: .status.phase}'",
                    "",
                    "# Pods con restart count > 5",
                    "kubectl get pods -A -o json | jq '.items[] |",
                    "  select(.status.containerStatuses[]?.restartCount > 5)'",
                ]),
                ("section", "Escenarios Complejos", ""),
                ("content", "App Intermitente", [
                    "1. ¿Pods reiniciando? → revisar restartCount",
                    "2. ¿Limits de recursos? → kubectl top pods",
                    "3. ¿OOMKilled? → describe pod → State",
                    "4. ¿Readiness probe intermitente? → describe pod → Conditions",
                ]),
                ("content", "Nodo NotReady", [
                    "1. Condiciones del nodo: kubectl describe node → Conditions",
                    "2. Status de kubelet: systemctl status kubelet",
                    "3. ¿Disco lleno? → Allocated resources section",
                    "4. Eventos del nodo: kubectl get events para el nodo",
                ]),
                ("code", "Checklist Rápido de Troubleshooting", [
                    "kubectl get pods -A              # Estado de todos los pods",
                    "kubectl get nodes -o wide        # Estado de todos los nodos",
                    "kubectl get events -A --sort-by='.lastTimestamp'",
                    "kubectl describe pod <nombre>    # ¡Events!",
                    "kubectl logs <pod> --previous    # Logs del crash",
                    "kubectl get svc,endpoints        # ¿Service conectado?",
                    "kubectl get netpol               # ¿Policies bloqueando?",
                    "az aks show --query provisioningState  # ¿Cluster OK?",
                ]),
                ("table", "Resumen del Curso", ["Lección", "Herramienta Clave"],
                 [["01 kubectl", "get, describe, logs"],
                  ["02 Pods", "describe pod → Events"],
                  ["03 Deployments", "rollout status, rollout undo"],
                  ["04 Services", "get endpoints, labels matching"],
                  ["05 ConfigMaps", "describe configmap/secret"],
                  ["06 Storage", "describe pvc, StorageClass"],
                  ["07 NetPol", "get netpol, selector matching"],
                  ["08 Nodos", "taints, labels, describe node"],
                  ["09 Azure", "az network nsg rule list"],
                  ["10 Avanzado", "Metodología DISCOVER"]]),
            ],
            "lab_title": "Desafío Multi-Problema",
            "lab_desc": "Un cluster tiene múltiples problemas combinados. Aplica todo lo aprendido.",
            "lab_cmd": "lab-10.sh",
        },
    },
}


def generate_lesson_pptx(num, lang, data):
    """Generate a single lesson PPTX file."""
    prs = create_presentation()

    # Title slide
    add_title_slide(prs, data["title"], data["subtitle"], num)

    # Content slides
    for slide_def in data["slides"]:
        stype = slide_def[0]
        if stype == "section":
            add_section_slide(prs, slide_def[1], slide_def[2] if len(slide_def) > 2 else "")
        elif stype == "content":
            add_content_slide(prs, slide_def[1], slide_def[2])
        elif stype == "code":
            add_code_slide(prs, slide_def[1], slide_def[2])
        elif stype == "table":
            add_table_slide(prs, slide_def[1], slide_def[2], slide_def[3])

    # Lab slide
    add_lab_slide(prs, num, data["lab_title"], data["lab_desc"], data["lab_cmd"], lang)

    # Save
    folder_names = {
        1: "01-fundamentos-kubectl", 2: "02-pods-y-contenedores",
        3: "03-deployments-y-replicas", 4: "04-servicios-y-networking",
        5: "05-configmaps-y-secrets", 6: "06-storage-y-volumenes",
        7: "07-network-policies", 8: "08-nodos-y-scheduling",
        9: "09-azure-networking", 10: "10-troubleshooting-avanzado",
    }
    folder = os.path.join(LESSONS_DIR, folder_names[num])
    os.makedirs(folder, exist_ok=True)

    suffix = "en" if lang == "en" else "es"
    filename = f"lesson-{num:02d}-{suffix}.pptx" if lang == "en" else f"leccion-{num:02d}-{suffix}.pptx"
    filepath = os.path.join(folder, filename)
    prs.save(filepath)
    return filepath


def main():
    print("Generating AKS Course Presentations...")
    print("=" * 50)

    for num in sorted(LESSONS.keys()):
        for lang in ["en", "es"]:
            data = LESSONS[num][lang]
            path = generate_lesson_pptx(num, lang, data)
            label = "EN" if lang == "en" else "ES"
            print(f"  [{label}] Lesson {num:02d}: {os.path.basename(path)}")

    print("=" * 50)
    print(f"Done! Generated {len(LESSONS) * 2} presentations.")


if __name__ == "__main__":
    main()
